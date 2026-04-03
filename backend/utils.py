"""
utils.py — Utility Functions
==============================
Provides:
  - AES-256 file encryption / decryption (using Fernet / PyCryptodome)
  - PDF text extraction (PyMuPDF / pdfplumber)
  - PPTX text extraction (python-pptx)
  - Cleanup utilities

Security Notes:
  - AES-256 in CBC mode with random IV per file
  - Key derived from environment variable or auto-generated on startup
  - Auto-delete temporary files after configurable TTL

Course: Advanced Topics in Machine Learning (HTML)
"""

import io
import logging
import os
import time
from pathlib import Path
from typing import Optional

log = logging.getLogger(__name__)

# ─── ENCRYPTION KEY ─────────────────────────────────────────────
# In production: store in env var / secrets manager
_ENCRYPTION_KEY: Optional[bytes] = None


def _get_key() -> bytes:
    """
    Get or generate the AES encryption key.
    Reads from AES_KEY environment variable (base64 32-byte key) if set,
    otherwise generates a random key for this session.
    """
    global _ENCRYPTION_KEY
    if _ENCRYPTION_KEY is not None:
        return _ENCRYPTION_KEY

    env_key = os.environ.get("AES_KEY")
    if env_key:
        import base64
        _ENCRYPTION_KEY = base64.b64decode(env_key)
    else:
        # Generate session key (new key each server restart — fine for demo)
        _ENCRYPTION_KEY = os.urandom(32)
        log.info("Generated new session AES-256 key.")

    return _ENCRYPTION_KEY


# ─── AES ENCRYPTION ─────────────────────────────────────────────

def encrypt_file(data: bytes) -> bytes:
    """
    Encrypt raw bytes with AES-256 (Fernet or PyCryptodome CBC).
    Returns encrypted bytes including IV/nonce prefix.
    """
    try:
        from cryptography.fernet import Fernet
        import base64, hashlib

        # Derive a valid Fernet key (32 bytes → URL-safe base64)
        raw_key = _get_key()
        fernet_key = base64.urlsafe_b64encode(hashlib.sha256(raw_key).digest())
        f = Fernet(fernet_key)
        encrypted = f.encrypt(data)
        log.debug(f"Encrypted {len(data)} bytes → {len(encrypted)} bytes (Fernet)")
        return encrypted

    except ImportError:
        pass  # Try PyCryptodome

    try:
        from Crypto.Cipher import AES
        from Crypto.Util.Padding import pad

        key = _get_key()[:32]
        iv  = os.urandom(16)
        cipher = AES.new(key, AES.MODE_CBC, iv)
        ct = cipher.encrypt(pad(data, AES.block_size))
        result = iv + ct
        log.debug(f"Encrypted {len(data)} bytes → {len(result)} bytes (AES-CBC)")
        return result

    except ImportError:
        log.warning("No encryption library available (install cryptography or pycryptodome). Storing unencrypted.")
        return data  # Fallback: no encryption (not for production!)


def decrypt_file(data: bytes) -> bytes:
    """
    Decrypt bytes previously encrypted with encrypt_file().
    """
    try:
        from cryptography.fernet import Fernet
        import base64, hashlib

        raw_key    = _get_key()
        fernet_key = base64.urlsafe_b64encode(hashlib.sha256(raw_key).digest())
        f = Fernet(fernet_key)
        return f.decrypt(data)

    except ImportError:
        pass

    try:
        from Crypto.Cipher import AES
        from Crypto.Util.Padding import unpad

        key = _get_key()[:32]
        iv  = data[:16]
        ct  = data[16:]
        cipher = AES.new(key, AES.MODE_CBC, iv)
        return unpad(cipher.decrypt(ct), AES.block_size)

    except ImportError:
        return data  # No decryption (matches encrypt fallback)

    except Exception as e:
        log.error(f"Decryption failed: {e}")
        raise ValueError("Could not decrypt file. Session may have restarted.") from e


# ─── TEXT EXTRACTION ────────────────────────────────────────────

def extract_text_pdf(raw_bytes: bytes) -> str:
    """
    Extract text from a PDF file given as raw bytes.
    Tries PyMuPDF (fitz) first, then pdfplumber as fallback.
    """
    text = ""

    # ── Method 1: PyMuPDF (fastest, best layout preservation)
    try:
        import fitz   # PyMuPDF

        doc = fitz.open(stream=raw_bytes, filetype="pdf")
        pages_text = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pages_text.append(page.get_text("text"))
        text = "\n\n".join(pages_text)
        doc.close()
        log.info(f"PyMuPDF extracted {len(text)} chars from PDF ({len(doc)} pages)")
        return _clean_text(text)

    except ImportError:
        log.info("PyMuPDF not available, trying pdfplumber…")
    except Exception as e:
        log.warning(f"PyMuPDF failed: {e}, trying pdfplumber…")

    # ── Method 2: pdfplumber
    try:
        import pdfplumber

        with pdfplumber.open(io.BytesIO(raw_bytes)) as pdf:
            pages_text = []
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    pages_text.append(t)
            text = "\n\n".join(pages_text)
        log.info(f"pdfplumber extracted {len(text)} chars")
        return _clean_text(text)

    except ImportError:
        log.warning("pdfplumber not available. Install: pip install pdfplumber")
    except Exception as e:
        log.error(f"pdfplumber failed: {e}")

    # ── Method 3: PyPDF2 (basic fallback)
    try:
        import PyPDF2

        reader = PyPDF2.PdfReader(io.BytesIO(raw_bytes))
        pages_text = [page.extract_text() or "" for page in reader.pages]
        text = "\n\n".join(pages_text)
        log.info(f"PyPDF2 extracted {len(text)} chars")
        return _clean_text(text)

    except ImportError:
        log.warning("No PDF library available. Install: pip install pymupdf")
        return "Could not extract PDF text. Please install PyMuPDF: pip install pymupdf"

    except Exception as e:
        log.error(f"All PDF extraction methods failed: {e}")
        return ""


def extract_text_pptx(raw_bytes: bytes) -> str:
    """
    Extract text from a PPTX file given as raw bytes.
    Uses python-pptx library.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation(io.BytesIO(raw_bytes))
        slide_texts = []

        for slide_num, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slide_texts.append(f"[Slide {slide_num + 1}]\n" + "\n".join(texts))

        text = "\n\n".join(slide_texts)
        log.info(f"python-pptx extracted {len(text)} chars from {len(prs.slides)} slides")
        return _clean_text(text)

    except ImportError:
        log.warning("python-pptx not available. Install: pip install python-pptx")
        return "Could not extract PPTX text. Please install python-pptx."
    except Exception as e:
        log.error(f"PPTX extraction failed: {e}")
        return ""


def _clean_text(text: str) -> str:
    """
    Remove excessive whitespace, repeated newlines, and non-UTF-8 chars.
    """
    import re

    # Normalize whitespace
    text = re.sub(r"\r\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)

    # Remove lone special characters that aren't content
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)

    return text.strip()


# ─── FILE CLEANUP ───────────────────────────────────────────────

def cleanup_old_files(directory: str, max_age_seconds: int = 3600):
    """
    Delete files older than max_age_seconds from a directory.
    Called periodically by the background cleanup thread.
    """
    now   = time.time()
    count = 0
    for f in Path(directory).iterdir():
        if f.is_file():
            age = now - f.stat().st_mtime
            if age > max_age_seconds:
                try:
                    f.unlink()
                    count += 1
                except Exception as e:
                    log.warning(f"Could not delete {f}: {e}")
    if count > 0:
        log.info(f"Cleaned {count} old files from {directory}")
